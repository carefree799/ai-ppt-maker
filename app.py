import os
import re
import json
import uuid
import requests
from flask import Flask, render_template, request, session, redirect, url_for, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
app.config['OUTPUT_FOLDER'] = os.path.join(os.path.dirname(__file__), 'output')
app.config['SECRET_KEY'] = os.urandom(24)
app.config['SESSION_DATA_DIR'] = '/tmp/ai-ppt-maker'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024

os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
os.makedirs(app.config['SESSION_DATA_DIR'], exist_ok=True)


def load_api_key():
    settings_path = os.path.join(os.path.dirname(__file__), 'settings.json')
    with open(settings_path, 'r') as f:
        settings = json.load(f)
    return settings.get('ANTHROPIC_API_KEY', '')


def generate_outline(topic, description=''):
    api_key = load_api_key()
    if not api_key:
        raise ValueError('API Key 未配置，请在 settings.json 中设置 ANTHROPIC_API_KEY')

    prompt = f"""你是一个 PPT 结构助手。根据用户的主题生成一个 PPT 大纲。

主题：{topic}
补充说明：{description}

请以 JSON 格式返回大纲，必须包含以下字段：
- title: PPT 标题
- slides: 数组，每页包含 page(页码，从1开始) 和 title(标题)

要求：
- 5-10 页为宜
- 逻辑清晰，符合主题
- 用中文回复
- 只返回 JSON，不要任何解释"""

    try:
        response = requests.post(
            'https://api.minimaxi.com/anthropic/v1/messages',
            headers={
                'Authorization': f'Bearer {api_key}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'MiniMax-M2.1',
                'max_tokens': 2000,
                'messages': [{'role': 'user', 'content': prompt}]
            },
            timeout=60
        )
    except requests.exceptions.Timeout:
        raise ValueError('API 请求超时，请重试')
    except requests.exceptions.RequestException as e:
        raise ValueError(f'API 请求失败: {str(e)}')

    if response.status_code != 200:
        try:
            error_msg = response.json().get('error', {}).get('message', '') or response.text
        except:
            error_msg = response.text
        raise ValueError(f'API 错误 ({response.status_code}): {error_msg[:100]}')

    try:
        result = response.json()
    except json.JSONDecodeError:
        raise ValueError('API 返回非 JSON 格式')

    content = result.get('content', [{}])[0].get('text', '')

    json_match = re.search(r'\{[\s\S]*\}', content)
    if json_match:
        try:
            outline = json.loads(json_match.group())

            if 'slides' in outline and isinstance(outline['slides'], list):
                for i, slide in enumerate(outline['slides']):
                    slide['page'] = i + 1

            return outline
        except json.JSONDecodeError:
            raise ValueError(f'无法解析 AI 输出的 JSON')
    raise ValueError(f'无法解析 AI 输出: {content[:200]}...' if len(content) > 200 else f'无法解析 AI 输出: {content}')


def generate_slides_content(session_data):
    api_key = load_api_key()
    topic = session_data['topic']
    description = session_data.get('description', '')
    outline = session_data['outline']

    prompt = f"""你是一个 PPT 内容助手。根据大纲为每一页生成详细内容。

主题：{topic}
补充说明：{description}

大纲：
{json.dumps(outline, ensure_ascii=False)}

请为每一页生成以下内容（返回 JSON 数组）：
- title: 标题
- bullets: 要点数组（3-5个）
- notes: 演讲备注（可选）

用中文回复，只返回 JSON 格式，不要任何解释"""

    try:
        response = requests.post(
            'https://api.minimaxi.com/anthropic/v1/messages',
            headers={
                'Authorization': f'Bearer {api_key}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'MiniMax-M2.1',
                'max_tokens': 4000,
                'messages': [{'role': 'user', 'content': prompt}]
            },
            timeout=90
        )
    except requests.exceptions.Timeout:
        raise ValueError('API 请求超时，请重试')
    except requests.exceptions.RequestException as e:
        raise ValueError(f'API 请求失败: {str(e)}')

    if response.status_code != 200:
        try:
            error_msg = response.json().get('error', {}).get('message', '') or response.text
        except:
            error_msg = response.text
        raise ValueError(f'API 错误 ({response.status_code}): {error_msg[:100]}')

    try:
        result = response.json()
    except json.JSONDecodeError:
        raise ValueError('API 返回非 JSON 格式')

    content = result.get('content', [{}])[0].get('text', '')

    json_match = re.search(r'\[[\s\S]*\]', content)
    if json_match:
        try:
            return json.loads(json_match.group())
        except json.JSONDecodeError:
            raise ValueError(f'无法解析 AI 输出的 JSON')
    raise ValueError(f'无法解析 AI 输出: {content[:200]}...' if len(content) > 200 else f'无法解析 AI 输出: {content}')
    raise ValueError(f'无法解析 AI 输出: {content[:200]}...' if len(content) > 200 else f'无法解析 AI 输出: {content}')


def create_pptx(slides_data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes

        title_box = shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        content_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True

        bullets = slide_data.get('bullets', [])
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f'• {bullet}'
            p.font.size = Pt(24)
            p.space_before = Pt(12)

    prs.save(output_path)


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/generate', methods=['POST'])
def generate():
    topic = request.form.get('topic', '').strip()
    description = request.form.get('description', '').strip()

    if not topic:
        return render_template('index.html', error='请输入主题')

    try:
        outline = generate_outline(topic, description)
        session_id = str(uuid.uuid4())

        session_data = {
            'topic': topic,
            'description': description,
            'outline': outline
        }
        session_file = os.path.join(app.config['SESSION_DATA_DIR'], f'{session_id}.json')
        with open(session_file, 'w') as f:
            json.dump(session_data, f, ensure_ascii=False)

        return render_template('preview.html', outline=outline, topic=topic,
                         description=description, session_id=session_id)
    except Exception as e:
        return render_template('index.html', error=f'生成失败: {str(e)}', topic=topic, description=description)


@app.route('/create', methods=['POST'])
def create():
    session_id = request.form.get('session_id', '')

    if not session_id:
        return render_template('index.html', error='无效的会话')

    session_file = os.path.join(app.config['SESSION_DATA_DIR'], f'{session_id}.json')
    if not os.path.exists(session_file):
        return render_template('index.html', error='会话已过期，请重新输入')

    with open(session_file, 'r') as f:
        session_data = json.load(f)

    try:
        slides_content = generate_slides_content(session_data)

        if not isinstance(slides_content, list):
            raise ValueError('生成的幻灯片内容格式不正确')

        output_file = f'{session_id}.pptx'
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_file)
        create_pptx(slides_content, output_path)

        return render_template('download.html', filename=output_file)
    except Exception as e:
        return render_template('index.html', error=f'生成失败: {str(e)}')


@app.route('/download/<path:filename>')
def download(filename):
    return send_from_directory(app.config['OUTPUT_FOLDER'], filename, as_attachment=True)


if __name__ == '__main__':
    app.run(debug=True)