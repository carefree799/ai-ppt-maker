"""
PPT Renderer - 使用 MCP 服务器生成专业级演示文稿
"""
import json
import os


class PPTRenderer:
    def __init__(self, template_path=None):
        self.template_path = template_path or os.environ.get('PPT_TEMPLATE_PATH', '/home/hly/templates')

    def render(self, slides_data, output_path):
        """
        渲染 PPT

        Args:
            slides_data: 幻灯片数据列表，每个元素包含：
                - title: 标题
                - bullets: 要点列表
                - layout: 布局类型
                - theme: 主题
                - chart_type: 图表类型(如需要)
                - chart_data: 图表数据(如需要)
            output_path: 输出文件路径
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RgbColor

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 主题颜色映射
        THEME_COLORS = {
            'blue': RgbColor(0, 120,212),
            'green': RgbColor(16,124,16),
            'purple': RgbColor(135,100,184),
            'gray': RgbColor(96,94,92)
        }

        for i, slide_data in enumerate(slides_data):
            layout = slide_data.get('layout', 'single_column')
            theme_name = slide_data.get('theme', 'blue' if i == 0 else 'gray')
            primary_color = THEME_COLORS.get(theme_name, THEME_COLORS['blue'])

            if i == 0 or layout == 'title_slide':
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_title_slide(slide, slide_data.get('title', ''), primary_color)
            elif layout == 'two_column' or layout == 'comparison':
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_two_column_slide(slide, slide_data, primary_color)
            elif layout == 'chart':
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_chart_slide(slide, slide_data, primary_color)
            else:  # single_column
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_single_column_slide(slide, slide_data, primary_color)

        prs.save(output_path)
        return output_path

    def _add_title_slide(self, slide, title, color):
        """添加标题页内容"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        shapes = slide.shapes

        # 主标题
        title_box = shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        subtitle_box = shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = 'AI 自动生成'
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    def _add_single_column_slide(self, slide, slide_data, color):
        """添加单栏内容页"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        shapes = slide.shapes

        # 标题
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # 内容区
        content_box = shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        bullets = slide_data.get('bullets', [])
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f'• {bullet}'
            p.font.size = Pt(18)
            p.space_before = Pt(12)

    def _add_two_column_slide(self, slide, slide_data, color):
        """添加双栏内容页"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        shapes = slide.shapes

        # 标题（横跨两栏）
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        bullets = slide_data.get('bullets', [])

        # 左栏
        left_box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True

        left_items = bullets[:len(bullets)//2] if len(bullets) > 1 else bullets
        for i, bullet in enumerate(left_items):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.text = f'• {bullet}'
            p.font.size = Pt(16)
            p.space_before = Pt(10)

        # 右栏
        right_box = shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True

        right_items = bullets[len(bullets)//2:] if len(bullets) > 1 else []
        for i, bullet in enumerate(right_items):
            if i == 0:
                p = tf_right.paragraphs[0]
            else:
                p = tf_right.add_paragraph()
            p.text = f'• {bullet}'
            p.font.size = Pt(16)
            p.space_before = Pt(10)

    def _add_chart_slide(self, slide, slide_data, color):
        """添加图表页（显示为表格/占位符）"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        shapes = slide.shapes

        # 标题
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # 图表占位符提示
        chart_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = chart_box.text_frame
        p = tf.paragraphs[0]

        chart_data = slide_data.get('chart_data', {})
        labels = chart_data.get('labels', [])
        values = chart_data.get('values', [])

        if labels and values:
            p.text = f"图表数据: {', '.join(labels)} 对应 {values}"
        else:
            p.text = "图表占位 - 请用 PowerPoint 编辑"

        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER


def render_pptx(slides_data, output_path):
    """渲染 PPT 的便捷函数"""
    renderer = PPTRenderer()
    return renderer.render(slides_data, output_path)